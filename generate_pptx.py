"""Generate an editable PowerPoint for 'The Engine & the Chassis' presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# Brand colors
GREEN = RGBColor(0x00, 0x9C, 0x6D)
GREEN_DARK = RGBColor(0x00, 0x6B, 0x4D)
GREEN_LIGHT = RGBColor(0x8D, 0xC9, 0xAB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1A, 0x1A, 0x2E)
GRAY_BG = RGBColor(0xF8, 0xF9, 0xFA)
GRAY_MID = RGBColor(0xE9, 0xEC, 0xEF)
ORANGE = RGBColor(0xD9, 0x77, 0x06)
BLUE = RGBColor(0x25, 0x63, 0xEB)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)
RED = RGBColor(0xDC, 0x26, 0x26)

# Slide dimensions (widescreen 16:9)
W = Inches(13.33)
H = Inches(7.5)


def new_prs():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    return prs


def blank_layout(prs):
    return prs.slide_layouts[6]  # blank


def add_rect(slide, x, y, w, h, fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, text, x, y, w, h,
                font_size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT,
                font_name="Calibri", italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox


def add_paragraph(tf, text, font_size=14, bold=False, color=DARK,
                  align=PP_ALIGN.LEFT, font_name="Calibri", italic=False,
                  space_before=Pt(4), bullet=False):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = space_before
    if bullet:
        p.level = 1
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return p


def add_slide_number(slide, num, total, color=DARK):
    add_textbox(slide, f"{num} / {total}",
                W - Inches(1.2), H - Inches(0.45), Inches(1.0), Inches(0.3),
                font_size=9, color=RGBColor(0xAA, 0xAA, 0xAA), align=PP_ALIGN.RIGHT)


def card(slide, x, y, w, h, title, body_lines,
         border_color=GREEN, bg_alpha_hint=None, title_color=DARK,
         body_color=DARK, title_size=16, body_size=12):
    """Draw a bordered card with title and bullet lines."""
    # Background
    bg = add_rect(slide, x, y, w, h,
                  fill_color=RGBColor(0xF0, 0xFB, 0xF6),
                  line_color=border_color, line_width=Pt(1.5))
    # Title
    tx = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.12),
                                  w - Inches(0.3), Inches(0.4))
    tx.word_wrap = True
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(title_size)
    run.font.bold = True
    run.font.color.rgb = title_color
    run.font.name = "Calibri"
    # Body lines
    tx2 = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.55),
                                   w - Inches(0.3), h - Inches(0.65))
    tx2.word_wrap = True
    tf2 = tx2.text_frame
    tf2.word_wrap = True
    for i, line in enumerate(body_lines):
        p2 = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p2.space_before = Pt(2)
        run2 = p2.add_run()
        run2.text = line
        run2.font.size = Pt(body_size)
        run2.font.color.rgb = body_color
        run2.font.name = "Calibri"


# ─────────────────────────────────────────────
# Slide builders
# ─────────────────────────────────────────────

def slide_hero(prs):
    slide = prs.slides.add_slide(blank_layout(prs))
    # Green gradient background
    add_rect(slide, 0, 0, W, H, fill_color=GREEN_DARK)
    add_rect(slide, 0, 0, W, H * 0.6, fill_color=GREEN)

    add_textbox(slide, "The Engine & the Chassis",
                Inches(0.8), Inches(1.8), Inches(11.7), Inches(1.4),
                font_size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, "How We Talk to AI Today",
                Inches(0.8), Inches(3.3), Inches(11.7), Inches(0.7),
                font_size=28, bold=False, color=RGBColor(0xFF, 0xFF, 0xFF),
                align=PP_ALIGN.CENTER)
    add_textbox(slide, "LLM = Engine  ·  Agentic Tool = Chassis",
                Inches(0.8), Inches(4.1), Inches(11.7), Inches(0.5),
                font_size=16, color=RGBColor(0xCC, 0xEE, 0xDD),
                align=PP_ALIGN.CENTER)
    return slide


def slide_agenda(prs, slide_num, total):
    slide = prs.slides.add_slide(blank_layout(prs))
    add_rect(slide, 0, 0, W, H, fill_color=GRAY_BG)
    add_textbox(slide, "Agenda", Inches(0.6), Inches(0.3), Inches(12), Inches(0.6),
                font_size=32, bold=True, color=GREEN)

    items = [
        ("🚗", "The Standard Interface", "Agentic tools today", "2 minutes"),
        ("🧩", "Two Tools to Know", "MCP servers & Agent Skills", "2 minutes"),
        ("⌨️", "opencode in action", "Live demo", "1 minute"),
    ]
    col_w = Inches(3.8)
    for i, (icon, title, desc, dur) in enumerate(items):
        x = Inches(0.6) + i * Inches(4.2)
        y = Inches(1.4)
        card(slide, x, y, col_w, Inches(4.5), f"{icon}  {title}",
             [desc, "", dur], border_color=GREEN_LIGHT, title_size=18,
             body_size=13, title_color=GREEN)

    add_slide_number(slide, slide_num, total)
    return slide


def slide_section(prs, part_num, part_title, subtitle, slide_num, total):
    slide = prs.slides.add_slide(blank_layout(prs))
    add_rect(slide, 0, 0, W, H, fill_color=GREEN_DARK)
    add_textbox(slide, f"Part {part_num}", Inches(0.8), Inches(1.8), Inches(11.7), Inches(0.7),
                font_size=20, color=GREEN_LIGHT, align=PP_ALIGN.CENTER)
    add_textbox(slide, part_title, Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.2),
                font_size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, subtitle, Inches(0.8), Inches(3.8), Inches(11.7), Inches(0.7),
                font_size=18, color=GREEN_LIGHT, align=PP_ALIGN.CENTER)
    add_slide_number(slide, slide_num, total)
    return slide


def slide_analogy(prs, slide_num, total):
    slide = prs.slides.add_slide(blank_layout(prs))
    add_rect(slide, 0, 0, W, H, fill_color=GRAY_BG)
    add_textbox(slide, "The Engine & the Chassis",
                Inches(0.6), Inches(0.25), Inches(12), Inches(0.55),
                font_size=28, bold=True, color=GREEN)

    # Left card — Engine
    card(slide, Inches(0.5), Inches(1.05), Inches(5.9), Inches(5.3),
         "🔧  Engine = the LLM",
         [
             "The model itself — raw intelligence and",
             "language understanding. It generates",
             "answers, writes code, reasons, translates.",
             "",
             "Examples: Claude 4, GPT-4o,",
             "Mistral Large, Llama 3…",
             "",
             "(italic) You don't drive a bare engine.",
         ],
         border_color=GREEN, title_color=GREEN, body_size=13)

    # Right card — Chassis
    card(slide, Inches(6.9), Inches(1.05), Inches(5.9), Inches(5.3),
         "🚗  Chassis = the Agentic Tool",
         [
             "The layer that makes the engine usable —",
             "connects it to files, tools, context,",
             "memory, and your workflow.",
             "",
             "It reads files, runs commands,",
             "browses the web, and iterates.",
             "",
             "(italic) It transforms raw power into a",
             "drivable experience.",
         ],
         border_color=BLUE, title_color=BLUE, body_size=13)

    add_slide_number(slide, slide_num, total)
    return slide


def slide_landscape(prs, slide_num, total):
    slide = prs.slides.add_slide(blank_layout(prs))
    add_rect(slide, 0, 0, W, H, fill_color=GRAY_BG)
    add_textbox(slide, "Landscape Overview",
                Inches(0.6), Inches(0.25), Inches(12), Inches(0.55),
                font_size=28, bold=True, color=GREEN)

    tools = [
        ("🤖", "Claude Code\n→ Claude Cowork", "Anthropic",
         "The first agentic tool. Cowork born when non-devs started using it — packaged for teams & entities.",
         "🔄 negotiation in progress", ORANGE),
        ("🧠", "Codex", "OpenAI",
         "Unified CLI + open-source option — two strong bets from OpenAI.",
         "✗ not available", RGBColor(0x64, 0x74, 0x8B)),
        ("🌊", "Mistral Vibe", "Mistral",
         "Our strategic partner's agentic tool.",
         "✓ available", GREEN),
        ("🏢", "IBM Bob", "IBM",
         "Experimentation currently underway internally.",
         "🧪 experimentation in progress", ORANGE),
        ("⚙️", "opencode", "Open source",
         "Most popular open-source agentic tool. Huge community ecosystem.",
         "✓ available", GREEN),
        ("➕", "And many more…", "Community & Big Tech",
         "Gemini CLI (Google), Continue (VS Code), OpenClaude, Amp, Cline, Cursor, Windsurf…",
         "ecosystem growing fast", RGBColor(0x64, 0x74, 0x8B)),
    ]

    # 3x2 grid
    col_w = Inches(4.0)
    col_h = Inches(2.8)
    cols = 3
    for i, (icon, name, maker, note, badge, badge_color) in enumerate(tools):
        col = i % cols
        row = i // cols
        x = Inches(0.4) + col * Inches(4.35)
        y = Inches(1.0) + row * Inches(3.05)
        bg_color = RGBColor(0xF0, 0xFB, 0xF6) if badge_color == GREEN else (
            RGBColor(0xFF, 0xF8, 0xF0) if badge_color == ORANGE else RGBColor(0xF8, 0xF9, 0xFA))
        line_style = Pt(1.5)
        add_rect(slide, x, y, col_w, col_h,
                 fill_color=bg_color, line_color=badge_color, line_width=line_style)
        add_textbox(slide, f"{icon}  {name}", x + Inches(0.15), y + Inches(0.12),
                    col_w - Inches(0.3), Inches(0.55),
                    font_size=13, bold=True, color=DARK)
        add_textbox(slide, maker, x + Inches(0.15), y + Inches(0.68),
                    col_w - Inches(0.3), Inches(0.25),
                    font_size=10, color=RGBColor(0x88, 0x88, 0x88))
        add_textbox(slide, note, x + Inches(0.15), y + Inches(0.95),
                    col_w - Inches(0.3), Inches(1.3),
                    font_size=10, color=DARK)
        add_textbox(slide, badge, x + Inches(0.15), y + col_h - Inches(0.45),
                    col_w - Inches(0.3), Inches(0.35),
                    font_size=10, bold=True, color=badge_color)

    add_slide_number(slide, slide_num, total)
    return slide


def slide_limiting_factor(prs, slide_num, total):
    slide = prs.slides.add_slide(blank_layout(prs))
    add_rect(slide, 0, 0, W, H, fill_color=GRAY_BG)
    add_textbox(slide, "The Limiting Factor",
                Inches(0.6), Inches(0.25), Inches(12), Inches(0.55),
                font_size=28, bold=True, color=GREEN)

    # Warning banner
    warn = add_rect(slide, Inches(0.6), Inches(1.1), Inches(12.1), Inches(1.4),
                    fill_color=RGBColor(0xFF, 0xF3, 0xE0), line_color=ORANGE, line_width=Pt(2))
    add_textbox(slide, "⚠️  A great chassis can't compensate for a weak engine",
                Inches(0.8), Inches(1.2), Inches(11.7), Inches(0.5),
                font_size=18, bold=True, color=ORANGE)
    add_textbox(slide,
                "The underlying LLM remains the limiting factor. Whichever agentic tool you use, "
                "the quality of its output depends on the model it runs on — its reasoning, knowledge, "
                "and ability to follow instructions.",
                Inches(0.8), Inches(1.72), Inches(11.7), Inches(0.65),
                font_size=13, color=DARK)

    # Two columns
    card(slide, Inches(0.6), Inches(2.9), Inches(5.8), Inches(3.8),
         "🚗  Chassis are now nearly feature-equivalent",
         ["• All major tools read files, run commands, use MCP",
          "• Real differentiator: open source vs. closed source",
          "• And popularity — community, plugins, integrations",
          "• Some vendors use the chassis for post-training",
          "  (fine-tuning on real agentic tasks) — a growing edge"],
         border_color=GREEN, title_color=GREEN, body_size=13)

    card(slide, Inches(6.9), Inches(2.9), Inches(5.8), Inches(3.8),
         "⚠️  The engine gap: frontier vs. small models",
         ["• Frontier models — Claude Opus/Sonnet, GPT-5, Kimi:",
          "  deep reasoning, large context, best results",
          "• Small models — Gemma, Mistral Mini, GPT-OSS:",
          "  faster & cheaper, but limited on complex tasks",
          "• Choosing the right model matters as much",
          "  as the chassis"],
         border_color=ORANGE, title_color=ORANGE, body_size=13)

    add_slide_number(slide, slide_num, total)
    return slide


def slide_two_standards(prs, slide_num, total):
    slide = prs.slides.add_slide(blank_layout(prs))
    add_rect(slide, 0, 0, W, H, fill_color=GRAY_BG)
    add_textbox(slide, "Two Open Standards",
                Inches(0.6), Inches(0.25), Inches(12), Inches(0.55),
                font_size=28, bold=True, color=GREEN)

    # AAIF banner
    add_rect(slide, Inches(0.6), Inches(1.05), Inches(12.1), Inches(0.6),
             fill_color=RGBColor(0xF3, 0xF0, 0xFF), line_color=PURPLE, line_width=Pt(1))
    add_textbox(slide,
                "Born at Anthropic  →  Now open standards under Linux Foundation — Agentic AI Foundation (AAIF)  •  Works with any LLM tool",
                Inches(0.8), Inches(1.12), Inches(11.7), Inches(0.45),
                font_size=11, color=PURPLE)

    # MCP card
    card(slide, Inches(0.5), Inches(1.85), Inches(5.9), Inches(4.7),
         "🔌  MCP Servers",
         ["Model Context Protocol — a universal adapter",
          "that plugs AI into your tools and data:",
          "databases, SaaS platforms, APIs, browsers…",
          "",
          '"Instead of copy-pasting, AI talks',
          'directly to your tools."'],
         border_color=GREEN, title_color=GREEN, body_size=13)

    # Skills card
    card(slide, Inches(6.9), Inches(1.85), Inches(5.9), Inches(4.7),
         "📋  Agent Skills",
         ["A reusable instruction sheet that teaches",
          "the AI a repeatable task — packaged once,",
          "available to everyone in your team.",
          "",
          '"Write the recipe once.',
          'Anyone can cook."'],
         border_color=PURPLE, title_color=PURPLE, body_size=13)

    add_slide_number(slide, slide_num, total)
    return slide


def slide_mcp(prs, slide_num, total):
    slide = prs.slides.add_slide(blank_layout(prs))
    add_rect(slide, 0, 0, W, H, fill_color=GRAY_BG)
    add_textbox(slide, "MCP Servers — Concrete Examples",
                Inches(0.6), Inches(0.25), Inches(12), Inches(0.55),
                font_size=28, bold=True, color=GREEN)

    add_textbox(slide, "What it looks like in practice",
                Inches(0.6), Inches(1.0), Inches(6), Inches(0.4),
                font_size=16, bold=True, color=DARK)

    examples = [
        ("🗄️ Postgres", '"How many users were created last month?"'),
        ("📚 Context7", '"Migrate this app using the official Angular 20 docs"'),
        ("☁️ IBM Cloud", '"Deploy the latest version to our UAT account"'),
        ("📋 JIRA", '"Create a bug ticket for this issue"'),
        ("🌐 Chrome Dev Tools", '"Improve the design of this page"'),
    ]
    for i, (server, prompt) in enumerate(examples):
        y = Inches(1.55) + i * Inches(0.9)
        add_rect(slide, Inches(0.6), y, Inches(5.9), Inches(0.78),
                 fill_color=RGBColor(0xF0, 0xFB, 0xF6),
                 line_color=GREEN_LIGHT, line_width=Pt(1))
        add_textbox(slide, server, Inches(0.75), y + Inches(0.08),
                    Inches(1.4), Inches(0.35), font_size=12, bold=True, color=GREEN)
        add_textbox(slide, prompt, Inches(2.2), y + Inches(0.08),
                    Inches(4.1), Inches(0.6), font_size=12, color=DARK, italic=True)

    # Right column
    add_textbox(slide, "The USB port of AI Agents",
                Inches(7.1), Inches(1.0), Inches(5.6), Inches(0.4),
                font_size=16, bold=True, color=DARK)

    card(slide, Inches(7.1), Inches(1.55), Inches(5.6), Inches(5.0),
         "Ecosystem",
         ["🗄️ Databases   ☁️ Cloud   📊 SaaS",
          "📖 Documentation   🔧 Dev tools",
          "",
          "Hundreds of community servers available"],
         border_color=GREEN_LIGHT, title_color=GREEN, body_size=13)

    add_slide_number(slide, slide_num, total)
    return slide


def slide_skills(prs, slide_num, total):
    slide = prs.slides.add_slide(blank_layout(prs))
    add_rect(slide, 0, 0, W, H, fill_color=GRAY_BG)
    add_textbox(slide, "Agent Skills — Concrete Examples",
                Inches(0.6), Inches(0.25), Inches(12), Inches(0.55),
                font_size=28, bold=True, color=GREEN)

    add_textbox(slide, "What it looks like in practice",
                Inches(0.6), Inches(1.0), Inches(6), Inches(0.4),
                font_size=16, bold=True, color=DARK)

    examples = [
        ("📄", "Document creation",
         '"Generate a project summary presentation" → PowerPoint following '
         'company template, with standard cover page, logo, and formatting.'),
        ("🎨", "Brand voice",
         '"Write a client email about the delay" → Drafted in the company tone, '
         'with correct sign-off, disclaimers, and style.'),
        ("📐", "Application templating",
         '"Bootstrap a new microservice" → Generates the full project structure '
         'following company standards: architecture, naming, CI config, README.'),
    ]
    for i, (icon, name, desc) in enumerate(examples):
        y = Inches(1.55) + i * Inches(1.5)
        add_rect(slide, Inches(0.6), y, Inches(5.9), Inches(1.3),
                 fill_color=RGBColor(0xF7, 0xF3, 0xFF),
                 line_color=PURPLE, line_width=Pt(1))
        add_textbox(slide, icon, Inches(0.75), y + Inches(0.15),
                    Inches(0.5), Inches(0.5), font_size=20)
        add_textbox(slide, name, Inches(1.35), y + Inches(0.1),
                    Inches(4.9), Inches(0.35), font_size=13, bold=True, color=PURPLE)
        add_textbox(slide, desc, Inches(1.35), y + Inches(0.48),
                    Inches(4.9), Inches(0.75), font_size=11, color=DARK, italic=True)

    # Right column
    add_textbox(slide, "The key idea",
                Inches(7.1), Inches(1.0), Inches(5.6), Inches(0.4),
                font_size=16, bold=True, color=DARK)

    add_rect(slide, Inches(7.1), Inches(1.55), Inches(5.6), Inches(1.1),
             fill_color=RGBColor(0xF7, 0xF3, 0xFF), line_color=PURPLE, line_width=Pt(1.5))
    add_textbox(slide,
                "📋  Write the recipe once.\nAnyone can cook — no technical skills needed.",
                Inches(7.25), Inches(1.65), Inches(5.3), Inches(0.9),
                font_size=13, color=DARK)

    card(slide, Inches(7.1), Inches(2.85), Inches(5.6), Inches(3.7),
         "Skills are shareable",
         ["A skill built by one team can be packaged",
          "and reused across the organization —",
          "or published to the community",
          "skills marketplace.",
          "",
          "No technical skills needed to use them."],
         border_color=PURPLE, title_color=PURPLE, body_size=13)

    add_slide_number(slide, slide_num, total)
    return slide


def slide_demo_template(prs, slide_num, total, title="opencode — Demo", subtitle="Add your demo scenario title here"):
    slide = prs.slides.add_slide(blank_layout(prs))
    add_rect(slide, 0, 0, W, H, fill_color=GRAY_BG)
    add_textbox(slide, title, Inches(0.6), Inches(0.25), Inches(12), Inches(0.55),
                font_size=28, bold=True, color=GREEN)

    # Badge
    add_rect(slide, Inches(0.6), Inches(1.0), Inches(1.8), Inches(0.38),
             fill_color=RGBColor(0xF0, 0xFB, 0xF6), line_color=GREEN, line_width=Pt(1))
    add_textbox(slide, "⌨️  opencode", Inches(0.65), Inches(1.05), Inches(1.7), Inches(0.3),
                font_size=11, bold=True, color=GREEN)

    # Subtitle placeholder
    add_textbox(slide, subtitle, Inches(2.6), Inches(1.05), Inches(10), Inches(0.3),
                font_size=12, color=RGBColor(0x99, 0x99, 0x99), italic=True)

    # Dashed placeholder box
    placeholder = add_rect(slide, Inches(0.6), Inches(1.6), Inches(12.1), Inches(5.3),
                            line_color=GREEN_LIGHT, line_width=Pt(1.5))
    placeholder.fill.background()

    add_textbox(slide, "🎬", Inches(5.8), Inches(2.8), Inches(1.5), Inches(0.7),
                font_size=36, align=PP_ALIGN.CENTER,
                color=RGBColor(0xCC, 0xCC, 0xCC))
    add_textbox(slide, "Demo content — replace with screenshots or steps",
                Inches(3.0), Inches(3.6), Inches(7.3), Inches(0.5),
                font_size=14, color=RGBColor(0xBB, 0xBB, 0xBB), align=PP_ALIGN.CENTER)
    add_textbox(slide, "Add screenshots, bullet steps, or a live walkthrough",
                Inches(3.0), Inches(4.1), Inches(7.3), Inches(0.4),
                font_size=12, color=RGBColor(0xCC, 0xCC, 0xCC),
                align=PP_ALIGN.CENTER, italic=True)

    add_slide_number(slide, slide_num, total)
    return slide


def slide_questions(prs, slide_num, total):
    slide = prs.slides.add_slide(blank_layout(prs))
    add_rect(slide, 0, 0, W, H, fill_color=GREEN_DARK)
    add_textbox(slide, "Questions?",
                Inches(0.8), Inches(1.8), Inches(11.7), Inches(1.2),
                font_size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, "This presentation was built WITH Claude Code",
                Inches(0.8), Inches(3.2), Inches(11.7), Inches(0.5),
                font_size=16, color=GREEN_LIGHT, align=PP_ALIGN.CENTER)

    contacts = [("📧", "Email"), ("💬", "Chat"), ("🔗", "Links")]
    for i, (icon, label) in enumerate(contacts):
        x = Inches(3.5) + i * Inches(2.2)
        add_textbox(slide, icon, x, Inches(4.2), Inches(2), Inches(0.6),
                    font_size=28, align=PP_ALIGN.CENTER, color=WHITE)
        add_textbox(slide, label, x, Inches(4.85), Inches(2), Inches(0.35),
                    font_size=13, align=PP_ALIGN.CENTER,
                    color=RGBColor(0xCC, 0xEE, 0xDD))

    add_slide_number(slide, slide_num, total)
    return slide


def slide_thankyou(prs, slide_num, total):
    slide = prs.slides.add_slide(blank_layout(prs))
    add_rect(slide, 0, 0, W, H, fill_color=GREEN_DARK)
    add_textbox(slide, "Thank You!",
                Inches(0.8), Inches(2.0), Inches(11.7), Inches(1.2),
                font_size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, "🔧  +  🚗  =  🚀",
                Inches(0.8), Inches(3.5), Inches(11.7), Inches(1.0),
                font_size=40, align=PP_ALIGN.CENTER, color=WHITE)
    add_slide_number(slide, slide_num, total)
    return slide


# ─────────────────────────────────────────────
# Build
# ─────────────────────────────────────────────

def build():
    prs = new_prs()
    total = 14

    slide_hero(prs)                                                          # 1
    slide_agenda(prs, 2, total)                                              # 2
    slide_section(prs, 1, "The Standard Interface",                          # 3
                  "Agentic tools are now the standard way to talk to an LLM",
                  3, total)
    slide_analogy(prs, 4, total)                                             # 4
    slide_landscape(prs, 5, total)                                           # 5
    slide_limiting_factor(prs, 6, total)                                     # 6
    slide_section(prs, 2, "Two Concepts to Know",                             # 7
                  "MCP Servers  ·  Agent Skills", 7, total)
    slide_two_standards(prs, 8, total)                                       # 8
    slide_mcp(prs, 9, total)                                                 # 9
    slide_skills(prs, 10, total)                                             # 10
    slide_section(prs, 3, "opencode in action",                              # 11
                  "Live demo", 11, total)
    slide_demo_template(prs, 12, total,                                      # 12
                        "opencode — Demo",
                        "Add your first demo scenario title here")
    slide_demo_template(prs, 13, total,                                      # 13
                        "opencode — Second Example",
                        "Add your second demo scenario title here")
    slide_questions(prs, 14, total)                                          # 14

    out = "slides-simple-editable.pptx"
    prs.save(out)
    print(f"✓ Saved: {out}")


if __name__ == "__main__":
    build()
